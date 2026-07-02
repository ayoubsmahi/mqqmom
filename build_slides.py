#!/usr/bin/env python3
"""Generate the final presentation PPTX from project content."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).parent
OUTPUT = ROOT / "MOM_IoT_Final_Praesentation.pptx"
DIAGRAM_SVG = ROOT / "mom_iot_system_diagram.svg"
DIAGRAM_PNG = ROOT / "mom_iot_system_diagram.png"

# Theme — matches mom_iot_system_diagram.svg
BG = RGBColor(0x0B, 0x0F, 0x14)
TITLE = RGBColor(0xF8, 0xFA, 0xFC)
BODY = RGBColor(0xCB, 0xD5, 0xE1)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2 = RGBColor(0x06, 0xB6, 0xD4)
FOOTER = RGBColor(0x64, 0x74, 0x8B)


def convert_diagram() -> None:
    if DIAGRAM_PNG.exists():
        return
    import subprocess

    subprocess.run(
        ["rsvg-convert", "-w", "1400", "-o", str(DIAGRAM_PNG), str(DIAGRAM_SVG)],
        check=True,
    )


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, prs: Presentation) -> None:
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def add_footer(slide, prs: Presentation, text: str = "Hochschule Merseburg · Rechnetzprojekt · SS26") -> None:
    box = slide.shapes.add_textbox(
        Inches(0.5),
        prs.slide_height - Inches(0.45),
        prs.slide_width - Inches(1),
        Inches(0.35),
    )
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.size = Pt(11)
    p.font.color.rgb = FOOTER
    p.alignment = PP_ALIGN.RIGHT


def add_title_slide(prs: Presentation, title: str, subtitle: str, extra: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_accent_bar(slide, prs)

    tbox = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.5))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = TITLE
    p.alignment = PP_ALIGN.CENTER

    sbox = slide.shapes.add_textbox(Inches(0.8), Inches(3.5), Inches(11.5), Inches(1.2))
    stf = sbox.text_frame
    stf.text = subtitle
    sp = stf.paragraphs[0]
    sp.font.size = Pt(22)
    sp.font.color.rgb = ACCENT2
    sp.alignment = PP_ALIGN.CENTER

    if extra:
        ebox = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.8))
        etf = ebox.text_frame
        etf.text = extra
        ep = etf.paragraphs[0]
        ep.font.size = Pt(18)
        ep.font.color.rgb = MUTED
        ep.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs)


def add_diagram_slide(prs: Presentation, title: str, image_path: Path, caption: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_accent_bar(slide, prs)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.7))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE

    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.35), Inches(1.05), width=Inches(12.6))

    cbox = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12), Inches(0.6))
    ctf = cbox.text_frame
    ctf.text = caption
    cp = ctf.paragraphs[0]
    cp.font.size = Pt(14)
    cp.font.color.rgb = ACCENT2
    cp.alignment = PP_ALIGN.CENTER

    add_footer(slide, prs)


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    *,
    subtitle: str = "",
    code: str = "",
    image_path: Path | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_accent_bar(slide, prs)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.9))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE

    top = 1.25
    if subtitle:
        sbox = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12), Inches(0.5))
        stf = sbox.text_frame
        stf.text = subtitle
        sp = stf.paragraphs[0]
        sp.font.size = Pt(16)
        sp.font.color.rgb = ACCENT2
        top += 0.55

    content_height = 5.5 if not image_path else 2.2
    bbox = slide.shapes.add_textbox(Inches(0.75), Inches(top), Inches(11.5), Inches(content_height))
    btf = bbox.text_frame
    btf.word_wrap = True

    for i, bullet in enumerate(bullets):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = bullet
        para.level = 0
        para.font.size = Pt(18)
        para.font.color.rgb = BODY
        para.space_after = Pt(10)

    if code:
        ctop = top + content_height + 0.1
        cbox = slide.shapes.add_textbox(Inches(0.75), Inches(ctop), Inches(11.5), Inches(1.2))
        ctf = cbox.text_frame
        ctf.text = code
        cp = ctf.paragraphs[0]
        cp.font.size = Pt(14)
        cp.font.name = "Consolas"
        cp.font.color.rgb = ACCENT2

    if image_path and image_path.exists():
        slide.shapes.add_picture(
            str(image_path),
            Inches(0.4),
            Inches(top + 0.2),
            width=Inches(12.5),
        )

    add_footer(slide, prs)


def add_table_slide(prs: Presentation, title: str, headers: list[str], rows: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_accent_bar(slide, prs)

    tbox = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(0.9))
    tf = tbox.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE

    cols = len(headers)
    table_shape = slide.shapes.add_table(len(rows) + 1, cols, Inches(0.6), Inches(1.3), Inches(12.0), Inches(0.4 * (len(rows) + 2)))
    table = table_shape.table

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1E, 0x27, 0x33)
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.bold = True
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = ACCENT2

    for row_idx, row in enumerate(rows, start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x15, 0x1B, 0x24)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.color.rgb = BODY

    add_footer(slide, prs)


def build() -> Path:
    convert_diagram()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Message Oriented Middleware (MOM)",
        "Finale Projektpräsentation — IoT Luftqualitätssystem",
        "Ayoub Smahi · Rechnetzprojekt · SS26",
    )

    add_content_slide(
        prs,
        "Agenda",
        [
            "1. Rückblick: Zwischenstand 1 & 2",
            "2. Gesamtarchitektur des Systems",
            "3. Hardware, Firmware & Serial-Bridge",
            "4. MQTT-Broker, Datenbank, Dashboard",
            "5. MOM-Konzepte in der Praxis",
            "6. Live-Demo",
            "7. Wireshark: Paketvisualisierung",
            "8. Fazit & Ausblick",
        ],
    )

    add_content_slide(
        prs,
        "Rückblick: Zwischenstand 1 — MOM-Theorie",
        [
            "Was ist MOM?",
            "• Vermittler (Broker) zwischen Sendern und Empfängern",
            "• Asynchrone Kommunikation",
            "• Publish/Subscribe-Modell (Basis für MQTT)",
            "",
            "Warum MOM?",
            "• Entkopplung · Skalierbarkeit · Zuverlässigkeit",
            "• MQTT + Mosquitto im IoT-Einsatz",
        ],
    )

    add_content_slide(
        prs,
        "Rückblick: Zwischenstand 2 — Umsetzung",
        [
            "Geplant am Ende von Stand 1:",
            "1. Grafisches Dashboard",
            "2. Datenbank (Persistenz)",
            "3. Warnsystem (Alerting)",
            "4. Fernsteuerung / Aktorik (Lüfter)",
            "",
            "In Stand 2 umgesetzt:",
            "✅ Serial-Bridge · ✅ SQLite · ✅ Alerts · ✅ Streamlit · ✅ Relais + Lüfter",
            "",
            "Heute: Live-Demo + Wireshark-Paketanalyse",
        ],
    )

    add_diagram_slide(
        prs,
        "Gesamtarchitektur (End-to-End)",
        DIAGRAM_PNG,
        "Topics: sensors/air_quality · devices/fan/control",
    )

    add_table_slide(
        prs,
        "Hardware & Firmware",
        ["Komponente", "Funktion"],
        [
            ["Elegoo Uno R3", "ATmega328P, USB-Serial 9600 Baud"],
            ["MQ135", "CO₂, Alkohol, Benzol — analog Pin A0"],
            ["5V-Relais", "Digital Pin 8 → Lüfter ein/aus"],
            ["main.cpp", "Alle 2 s: analogRead → Serial"],
            ["main.cpp", "Empfängt ON/OFF → Relais schalten"],
        ],
    )

    add_content_slide(
        prs,
        "Serial-Bridge — Zwei-Wege-Kommunikation",
        [
            "Elegoo ←—— USB Serial ——→ mqtt_bridge.py ←—— MQTT ——→ Mosquitto",
            "",
            "Richtung 1 — Sensor → Netz:",
            "Bridge liest Zahlen von USB → publish sensors/air_quality",
            "",
            "Richtung 2 — Netz → Aktor:",
            "Bridge subscribed devices/fan/control → ON/OFF an Elegoo",
            "",
            "MOM-Vorteil: Sensor kennt Dashboard/DB nicht — lose Kopplung",
        ],
    )

    add_content_slide(
        prs,
        "MQTT-Broker (Mosquitto)",
        [
            "Broker: Mosquitto · localhost:1883",
            "",
            "Topic sensors/air_quality → Payload: 239, 312, …",
            "Topic devices/fan/control → Payload: ON, OFF",
            "",
            "Clients:",
            "• mqtt_bridge.py — publish + subscribe",
            "• db_logger_alerter.py — subscribe",
            "• dashboard.py — publish (Fan-Steuerung)",
        ],
    )

    add_content_slide(
        prs,
        "Datenbank & Warnsystem",
        [
            "db_logger_alerter.py — unabhängiger MQTT-Subscriber",
            "",
            "• Speichert alle Werte in SQLite (air_quality_data.db)",
            "• Tabelle sensor_logs: id · timestamp · value",
            "• Alert-Schwellwert 300 → CRITICAL WARNING im Terminal",
            "",
            "MOM-Vorteil: Nachträglich hinzugefügt — ohne Elegoo-Code zu ändern",
        ],
    )

    add_content_slide(
        prs,
        "Dashboard (Streamlit)",
        [
            "dashboard.py · http://localhost:8501",
            "",
            "• Live-Metrik — aktueller Luftqualitätswert",
            "• Liniendiagramm — letzte 50 Messungen",
            "• Kritische Warnung bei Wert > 300",
            "• Buttons: Turn Fan ON / OFF",
            "• Auto-Refresh alle 2 Sekunden",
            "",
            "Dashboard und Sensor kommunizieren nur über den Broker",
        ],
    )

    add_table_slide(
        prs,
        "MOM-Konzepte: Theorie trifft Praxis",
        ["MOM-Konzept", "Umsetzung im Projekt"],
        [
            ["Publish/Subscribe", "MQTT-Topics, Mosquitto"],
            ["Entkopplung", "C++ · Python · Web · Hardware"],
            ["Skalierbarkeit", "Logger + Dashboard ohne Elegoo-Änderung"],
            ["Asynchronität", "Sensor sendet, Logger speichert unabhängig"],
            ["Bidirektionalität", "Sensordaten raus, Lüfterbefehle rein"],
        ],
    )

    add_content_slide(
        prs,
        "Hinweis: Hardware-Simulator",
        [
            "⚠ Physisches Elegoo-Board derzeit nicht verfügbar",
            "✅ Software-Simulator mit identischem Protokoll",
            "✅ Gesamte Software-Pipeline unverändert",
            "",
            "Alle 2 s Sensorwert · ON/OFF für Relais",
            "Bridge · MQTT · DB · Dashboard · Wireshark — identisch",
            "",
            "Stärke der Architektur: Software unabhängig von Hardware",
        ],
    )

    add_content_slide(
        prs,
        "Live-Demo — System starten",
        [
            "Start mit einem Befehl:",
            "",
            "Läuft parallel:",
            "• Mosquitto Broker",
            "• mqtt_bridge.py (simuliert)",
            "• db_logger_alerter.py",
            "• dashboard.py (Streamlit)",
        ],
        code="./run_demo.sh",
    )

    add_content_slide(
        prs,
        "Live-Demo — Dashboard & Steuerung",
        [
            "Browser: http://localhost:8501",
            "",
            "Zeigen:",
            "1. Live-Graph — Werte alle 2 Sekunden",
            "2. Alert — Wert > 300 → kritische Warnung",
            "3. Turn Fan ON → MQTT publish devices/fan/control",
            "4. Bridge: Relay ON (simulated)",
            "",
            "Regelkreis: Messung → Speicherung → Visualisierung → Reaktion",
        ],
    )

    add_content_slide(
        prs,
        "Live-Demo — Wireshark",
        [
            "Interface: Loopback (lo)",
            "Filter: mqtt",
            "",
            "Zeigen:",
            "• PUBLISH sensors/air_quality — Payload z. B. 239",
            "• PUBLISH devices/fan/control — Payload ON / OFF",
            "",
            "Nachweis: Echte MQTT-Pakete auf TCP Port 1883",
            "CONNECT · SUBSCRIBE · PUBLISH — wie in produktiver IoT-Umgebung",
        ],
    )

    add_content_slide(
        prs,
        "Fazit",
        [
            "Vom Konzept zur fertigen Lösung:",
            "1. Stand 1: MOM-Theorie (Pub/Sub, Entkopplung)",
            "2. Stand 2: IoT-Implementierung (Bridge, DB, Dashboard, Aktorik)",
            "3. Finale: Live-Demo + Wireshark-Nachweis",
            "",
            "Ergebnis: Bidirektionales IoT-System mit MOM",
            "",
            "Lose Kopplung · Skalierbarkeit · Asynchrone Kommunikation",
            "C++ · Python · Web · 5V-Hardware · 12V-Aktorik — ein Broker",
        ],
    )

    add_content_slide(
        prs,
        "Quellen",
        [
            "• G2 — Message-Oriented Middleware",
            "• GeeksforGeeks — What is MOM",
            "• HPI — MOM Paper (mds04_mom.pdf)",
            "• Mosquitto / MQTT Dokumentation",
            "• Streamlit · PlatformIO · Wireshark",
        ],
    )

    add_title_slide(
        prs,
        "Vielen Dank für Ihre Aufmerksamkeit",
        "Fragen?",
        "MOM · MQTT · Serial-Bridge · Hardware · Wireshark",
    )

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
